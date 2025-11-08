"""
PowerPoint document processing module - テキスト、表、画像の抽出
"""

import logging
from pathlib import Path
from typing import Dict, List, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tiktoken
from PIL import Image
import io
from langchain_text_splitters import RecursiveCharacterTextSplitter


logger = logging.getLogger(__name__)


class PowerPointProcessor:
    """PowerPointファイルからテキスト、表、画像を抽出するクラス"""

    def __init__(self, config: Dict[str, Any]):
        """
        初期化

        Args:
            config: 処理設定
        """
        self.config = config
        self.pdf_config = config.get("pdf_processing", {})  # 同じ設定を使用
        self.vision_config = config.get("vision", {})
        self.rag_config = config.get("rag", {})
        self.chunking_config = config.get("chunking", {})

        # tiktokenエンコーダーの初期化
        model_name = config.get("openai", {}).get("model_chat", "gpt-4.1")
        try:
            self.encoding = tiktoken.encoding_for_model(model_name)
        except KeyError:
            self.encoding = tiktoken.get_encoding("cl100k_base")

        # セマンティックチャンカーの初期化
        self.text_splitter = None
        if self.rag_config.get("enable_semantic_chunking", False):
            try:
                chunk_size = self.chunking_config.get("chunk_size", self.pdf_config.get("chunk_size", 800))
                chunk_overlap = self.chunking_config.get("chunk_overlap", self.pdf_config.get("chunk_overlap", 150))
                separators = self.chunking_config.get("separators", ["\n\n", "\n", "。", "．", ". ", "! ", "? ", "；", "、", ", ", " ", ""])

                self.text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                    length_function=lambda text: len(self.encoding.encode(text)),
                    separators=separators,
                    keep_separator=True
                )
                logger.info(f"Semantic chunking initialized for PowerPoint (chunk_size={chunk_size}, overlap={chunk_overlap})")
            except Exception as e:
                logger.error(f"Failed to initialize semantic chunker: {e}")
                logger.warning("Falling back to legacy chunking")

    def process_powerpoint(self, pptx_path: str, category: str) -> Dict[str, Any]:
        """
        PowerPointファイルを処理してテキスト、表、画像を抽出

        Args:
            pptx_path: PowerPointファイルのパス
            category: ドキュメントカテゴリー

        Returns:
            dict: 抽出結果（テキストチャンク、画像情報など）
        """
        logger.info(f"Processing PowerPoint document: {pptx_path}")

        result = {
            "source_file": Path(pptx_path).name,
            "category": category,
            "text_chunks": [],
            "images": [],
            "metadata": {},
        }

        try:
            # PowerPointファイルを開く
            prs = Presentation(pptx_path)

            # メタデータ
            result["metadata"]["slides_count"] = len(prs.slides)

            # 画像出力ディレクトリを準備
            doc_name = Path(pptx_path).stem
            output_dir = Path(f"data/extracted_images/{doc_name}")
            output_dir.mkdir(parents=True, exist_ok=True)

            # 各スライドを処理
            image_index = 0

            for slide_num, slide in enumerate(prs.slides, start=1):
                # 10スライドごと、または最終スライドでログ出力
                if slide_num % 10 == 0 or slide_num == len(prs.slides):
                    logger.info(f"Processing slide {slide_num}/{len(prs.slides)}")

                slide_text = []

                # スライド番号をヘッダーとして追加
                slide_text.append(f"\n\n## スライド {slide_num}\n\n")

                # スライド内の各シェイプを処理
                for shape in slide.shapes:
                    # テキストを抽出
                    if hasattr(shape, "text") and shape.text.strip():
                        # タイトルの場合は見出しとして
                        if shape == slide.shapes.title:
                            slide_text.append(f"### {shape.text.strip()}\n\n")
                        else:
                            slide_text.append(f"{shape.text.strip()}\n\n")

                    # 表を抽出
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_markdown = self._table_to_markdown(shape.table)
                        if table_markdown:
                            slide_text.append(f"[表]\n{table_markdown}\n\n")

                    # 画像を抽出
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob

                            # PIL Imageに変換
                            pil_image = Image.open(io.BytesIO(image_bytes))

                            # 画像をリサイズ（大きすぎる場合）
                            max_size = self.vision_config.get("max_image_size", 2000)
                            if max(pil_image.size) > max_size:
                                ratio = max_size / max(pil_image.size)
                                new_size = tuple(int(dim * ratio) for dim in pil_image.size)
                                pil_image = pil_image.resize(new_size, Image.Resampling.LANCZOS)

                            # 画像を保存
                            image_filename = f"slide_{slide_num}_image_{image_index}.png"
                            image_path = output_dir / image_filename
                            pil_image.save(str(image_path), format=self.vision_config.get("image_format", "PNG"))

                            # メタデータを作成
                            image_info = {
                                "image_path": str(image_path),
                                "page_number": slide_num,  # スライド番号
                                "image_index": image_index,
                                "width": pil_image.size[0],
                                "height": pil_image.size[1],
                                "source_file": result["source_file"],
                                "category": category,
                                "content_type": "image",
                                "slide_number": slide_num
                            }
                            result["images"].append(image_info)
                            image_index += 1

                            logger.debug(f"Extracted image {image_index} from slide {slide_num}")

                        except Exception as e:
                            logger.warning(f"Failed to extract image from slide {slide_num}: {e}")
                            continue

                # スライドのテキストを結合
                full_slide_text = "".join(slide_text)

                # テキストをチャンクに分割
                if full_slide_text.strip():
                    chunks = self._create_text_chunks(
                        full_slide_text,
                        section_num=slide_num,  # スライド番号をセクション番号として使用
                        source_file=result["source_file"],
                        category=category,
                        slide_number=slide_num
                    )
                    result["text_chunks"].extend(chunks)

            logger.info(f"Extracted {len(result['text_chunks'])} text chunks and {len(result['images'])} images from PowerPoint document")

        except Exception as e:
            logger.error(f"Error processing PowerPoint document: {e}")
            raise

        return result

    def _table_to_markdown(self, table) -> str:
        """
        PowerPoint表をMarkdown形式に変換

        Args:
            table: python-pptxのTableオブジェクト

        Returns:
            str: Markdown形式の表
        """
        if not table.rows:
            return ""

        markdown_lines = []

        # ヘッダー行（最初の行）
        header_cells = [cell.text.strip() for cell in table.rows[0].cells]
        header = "| " + " | ".join(header_cells) + " |"
        markdown_lines.append(header)

        # セパレーター
        separator = "| " + " | ".join(["---"] * len(header_cells)) + " |"
        markdown_lines.append(separator)

        # データ行
        for row in table.rows[1:]:
            row_cells = [cell.text.strip() for cell in row.cells]
            data_row = "| " + " | ".join(row_cells) + " |"
            markdown_lines.append(data_row)

        return "\n".join(markdown_lines)

    def _create_text_chunks(
        self,
        text: str,
        section_num: int,
        source_file: str,
        category: str,
        slide_number: int = None
    ) -> List[Dict[str, Any]]:
        """
        テキストをチャンクに分割

        Args:
            text: テキスト
            section_num: セクション番号（スライド番号）
            source_file: ソースファイル名
            category: カテゴリー
            slide_number: スライド番号

        Returns:
            list: テキストチャンクのリスト
        """
        # セマンティックチャンキング or レガシーチャンキング
        if self.text_splitter:
            # セマンティックチャンキング使用
            chunk_texts = self.text_splitter.split_text(text)
            chunks = []
            for idx, chunk_text in enumerate(chunk_texts):
                if chunk_text.strip():
                    token_count = len(self.encoding.encode(chunk_text))
                    chunk_data = {
                        "text": chunk_text.strip(),
                        "page_number": section_num,  # スライド番号
                        "source_file": source_file,
                        "category": category,
                        "chunk_index": idx,
                        "token_count": token_count,
                    }
                    if slide_number:
                        chunk_data["slide_number"] = slide_number
                    chunks.append(chunk_data)
            logger.debug(f"Created {len(chunks)} semantic chunks from PowerPoint slide")
            return chunks
        else:
            # レガシーチャンキング
            return self._legacy_create_text_chunks(text, section_num, source_file, category, slide_number)

    def _legacy_create_text_chunks(
        self,
        text: str,
        section_num: int,
        source_file: str,
        category: str,
        slide_number: int = None
    ) -> List[Dict[str, Any]]:
        """
        テキストをチャンクに分割（レガシー実装）

        Args:
            text: テキスト
            section_num: セクション番号
            source_file: ソースファイル名
            category: カテゴリー
            slide_number: スライド番号

        Returns:
            list: テキストチャンクのリスト
        """
        chunk_size = self.pdf_config.get("chunk_size", 800)
        chunk_overlap = self.pdf_config.get("chunk_overlap", 150)
        chunks = []

        # テキストをトークン化
        tokens = self.encoding.encode(text)

        # トークン数が少ない場合はそのまま返す
        if len(tokens) <= chunk_size:
            chunk_data = {
                "text": text.strip(),
                "page_number": section_num,
                "source_file": source_file,
                "category": category,
                "chunk_index": 0,
                "token_count": len(tokens)
            }
            if slide_number:
                chunk_data["slide_number"] = slide_number
            return [chunk_data]

        # チャンクに分割
        start = 0
        chunk_index = 0

        while start < len(tokens):
            end = min(start + chunk_size, len(tokens))
            chunk_tokens = tokens[start:end]
            chunk_text = self.encoding.decode(chunk_tokens)

            if chunk_text.strip():
                chunk_data = {
                    "text": chunk_text.strip(),
                    "page_number": section_num,
                    "source_file": source_file,
                    "category": category,
                    "chunk_index": chunk_index,
                    "token_count": len(chunk_tokens),
                }
                if slide_number:
                    chunk_data["slide_number"] = slide_number
                chunks.append(chunk_data)
                chunk_index += 1

            start = end - chunk_overlap if end < len(tokens) else end

        logger.debug(f"Created {len(chunks)} chunks from PowerPoint slide ({len(tokens)} tokens)")
        return chunks
